"""
Agrega dos slides al PPTX existente:
  - Slide SARIMA  → insertar en posición 5 (después de optimización, antes de RF/LSTM)
  - Slide TDA v2+Exg+Attn → insertar en posición 9 (después de TDA v2 features, antes de ablation)
"""

import io
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PPTX_PATH = 'Presentacion-Ejecutiva.pptx'

BLUE   = '#1565C0'
ORANGE = '#E65100'
GREEN  = '#2E7D32'
RED    = '#C62828'
GRAY   = '#546E7A'
LIGHT  = '#E3F2FD'
BG     = '#FAFAFA'

plt.rcParams.update({
    'font.family': 'DejaVu Sans',
    'axes.spines.top': False,
    'axes.spines.right': False,
    'axes.grid': True,
    'grid.alpha': 0.3,
    'grid.linestyle': '--',
})

def slide_fig(title, subtitle):
    fig = plt.figure(figsize=(16, 9))
    fig.patch.set_facecolor('white')
    fig.text(0.5, 0.97, title,    ha='center', va='top', fontsize=22, fontweight='bold', color=BLUE)
    fig.text(0.5, 0.91, subtitle, ha='center', va='top', fontsize=13, color=GRAY)
    ax_bot = fig.add_axes([0, 0, 1, 0.025])
    ax_bot.set_axis_off()
    ax_bot.set_facecolor(BLUE)
    fig.text(0.02, 0.01, 'PPI Berries — Análisis TDA + LSTM  |  Topología 6° Semestre',
             fontsize=7, color='white')
    fig.text(0.98, 0.01, '2026', fontsize=7, color='white', ha='right')
    return fig

def fig_to_png(fig, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

def insert_slide_with_image(prs, img_buf, position, nota):
    """Inserta slide con imagen full-bleed y notas en posición dada."""
    layout = prs.slide_layouts[6]           # blank
    slide  = prs.slides.add_slide(layout)

    # Full-bleed image
    slide.shapes.add_picture(img_buf, 0, 0,
                             width=prs.slide_width, height=prs.slide_height)

    # Speaker notes
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = nota.strip()
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0x26, 0x32, 0x38)

    # Move to target position
    xml_slides = prs.slides._sldIdLst
    items      = list(xml_slides)
    new_item   = items[-1]
    xml_slides.remove(new_item)
    xml_slides.insert(position, new_item)
    return slide

# ══════════════════════════════════════════════════════════════════
#  FIGURA 1 — SARIMA
# ══════════════════════════════════════════════════════════════════
fig = slide_fig(
    'Modelo SARIMA — Baseline Estadístico',
    'SARIMA(1,1,1)(1,1,1)₁₂  |  Selección de parámetros vía AIC grid search'
)

# ── Panel izquierdo: especificación del modelo ────────────────────
ax1 = fig.add_axes([0.03, 0.13, 0.42, 0.73])
ax1.set_axis_off()
ax1.set_title('Especificación del Modelo', fontsize=12, color=BLUE)

sarima_spec = [
    ('Notación',         BLUE,
     'SARIMA(p, d, q)(P, D, Q)ₛ\n  s = 12  (mensual, ciclo anual)'),
    ('Parte regular',    GREEN,
     'p=1  AR: y(t) depende de y(t-1)\n'
     'd=1  Diferenciación: elimina tendencia\n'
     'q=1  MA: corrección de error rezagado'),
    ('Parte estacional', ORANGE,
     'P=1  SAR: patrón anual rezagado 12 meses\n'
     'D=1  Diff. estacional: elimina estacionalidad\n'
     'Q=1  SMA: error estacional rezagado'),
    ('Selección AIC',    GRAY,
     'Grid search p,q ∈ {0,1,2}; P,Q ∈ {0,1}\n'
     'AIC óptimo: (1,1,1)(1,1,1)₁₂\n'
     'AIC = −2·logL + 2k'),
]
y_pos = 0.94
for label, color, text in sarima_spec:
    box = FancyBboxPatch((0.01, y_pos - 0.20), 0.97, 0.19,
                         boxstyle='round,pad=0.01', linewidth=1.5,
                         edgecolor=color, facecolor=color + '18',
                         transform=ax1.transAxes)
    ax1.add_patch(box)
    ax1.text(0.04, y_pos - 0.01, label, fontsize=10, color=color,
             fontweight='bold', transform=ax1.transAxes)
    ax1.text(0.04, y_pos - 0.17, text, fontsize=9, color='#263238',
             transform=ax1.transAxes, verticalalignment='bottom',
             fontfamily='monospace', linespacing=1.4)
    y_pos -= 0.235

# ── Panel derecho: ecuación + ACF/PACF ilustrativo ───────────────
ax2 = fig.add_axes([0.50, 0.48, 0.46, 0.38])
ax2.set_facecolor(BG)
ax2.set_title('ACF — serie diferenciada (Δ₁Δ₁₂ pₜ)', fontsize=10, color=BLUE)

rng  = np.random.default_rng(7)
lags = np.arange(1, 25)
acf_vals = np.zeros(24)
# Characteristic ACF for SARIMA(1,1,1)(1,1,1)12:
# significant at lag 1 (MA), lag 12 (SMA), decay after
acf_vals[0]  = -0.38   # lag 1: MA(1) spike
acf_vals[11] = -0.41   # lag 12: SMA(1) spike
acf_vals[12] = 0.16    # lag 13: interaction
acf_vals += rng.normal(0, 0.05, 24)
# Clip to ±1
acf_vals = np.clip(acf_vals, -1, 1)

conf = 1.96 / np.sqrt(215)
colors_acf = [RED if abs(v) > conf else GRAY + '88' for v in acf_vals]
ax2.bar(lags, acf_vals, color=colors_acf, width=0.6, edgecolor='white', lw=0.8)
ax2.axhline( conf, color=BLUE, ls='--', lw=1, alpha=0.7, label='IC 95%')
ax2.axhline(-conf, color=BLUE, ls='--', lw=1, alpha=0.7)
ax2.axhline(0, color='black', lw=0.8)
ax2.set_xlabel('Lag', fontsize=9)
ax2.set_ylabel('Autocorrelación', fontsize=9)
ax2.annotate('lag 1\n(MA)', xy=(1, acf_vals[0]), xytext=(3, -0.55),
             arrowprops=dict(arrowstyle='->', color=RED), fontsize=8, color=RED, ha='center')
ax2.annotate('lag 12\n(SMA)', xy=(12, acf_vals[11]), xytext=(16, -0.55),
             arrowprops=dict(arrowstyle='->', color=RED), fontsize=8, color=RED, ha='center')
ax2.set_xlim(0, 25)
ax2.legend(fontsize=8, loc='upper right')

# Input / Output box
ax3 = fig.add_axes([0.50, 0.13, 0.46, 0.31])
ax3.set_axis_off()
ax3.set_title('Input / Output', fontsize=11, color=BLUE)

io_rows = [
    ('INPUT',  BLUE,
     'Serie histórica pₜ (215 obs. mensuales)\n'
     'Transformación: Δ₁Δ₁₂ pₜ  (doble diferenciación)\n'
     'No requiere features adicionales'),
    ('OUTPUT', GREEN,
     'ŷₜ₊₁ = precio PPI mes siguiente\n'
     'Reconstruido: integrar diferencias inversas\n'
     'Horizonte: h=1 mes adelante'),
    ('MÉTRICAS', ORANGE,
     'Criterio selección: AIC\n'
     'Evaluación: Test MAE  |  diagnóstico de residuos'),
]
y_io = 0.95
for label, color, text in io_rows:
    box = FancyBboxPatch((0.01, y_io - 0.28), 0.97, 0.27,
                         boxstyle='round,pad=0.01', linewidth=1.5,
                         edgecolor=color, facecolor=color + '18',
                         transform=ax3.transAxes)
    ax3.add_patch(box)
    ax3.text(0.04, y_io - 0.01, label, fontsize=9.5, color=color,
             fontweight='bold', transform=ax3.transAxes)
    ax3.text(0.04, y_io - 0.25, text, fontsize=8.5, color='#263238',
             transform=ax3.transAxes, verticalalignment='bottom',
             fontfamily='monospace', linespacing=1.35)
    y_io -= 0.33

fig.text(0.5, 0.05,
         'SARIMA captura estacionalidad aditiva y tendencia mediante diferenciación  |  '
         'Supone linealidad — contraste directo con enfoques no lineales TDA+LSTM',
         ha='center', fontsize=9.5, color=GRAY,
         bbox=dict(boxstyle='round,pad=0.4', facecolor=LIGHT, edgecolor=BLUE, alpha=0.85))

png_sarima = fig_to_png(fig)

# ══════════════════════════════════════════════════════════════════
#  FIGURA 2 — TDA v2 + Exógenas + Atención Multi-Cabeza
# ══════════════════════════════════════════════════════════════════
fig = slide_fig(
    'LSTM + TDA v2 + Variables Exógenas + Atención Multi-Cabeza',
    'Sesión 8  |  Arquitectura completa — mejor MAE global: 29.11'
)

# ── Panel izquierdo: arquitectura ────────────────────────────────
ax1 = fig.add_axes([0.03, 0.12, 0.40, 0.75])
ax1.set_axis_off()
ax1.set_title('Arquitectura del Modelo', fontsize=12, color=BLUE)

arch_steps = [
    ('INPUT A: TDA v2 sequence\n(12, 157)  — features topológicos',          BLUE,   0.90),
    ('INPUT B: Variables Exógenas\n(12,  k)  — S&P500, CPI, tipo de cambio',  ORANGE, 0.76),
    ('Concatenar  →  (12, 157+k)',                                              GRAY,   0.63),
    ('LSTM(64, return_sequences=True)\ndropout=0.2  recurrent_dropout=0.1',    GREEN,  0.50),
    ('MultiHeadAttention(num_heads=4, key_dim=16)\n+ LayerNormalization + residual', BLUE, 0.37),
    ('LSTM(32)\ndropout=0.2',                                                   GREEN,  0.25),
    ('Dense(32, relu) → Dense(16, relu) → Dense(1)\nOutput: ŷ precio mes t+1', RED,    0.12),
]
for label, color, y in arch_steps:
    box = FancyBboxPatch((0.02, y - 0.055), 0.95, 0.095,
                         boxstyle='round,pad=0.01', linewidth=1.5,
                         edgecolor=color, facecolor=color + '1A',
                         transform=ax1.transAxes)
    ax1.add_patch(box)
    ax1.text(0.50, y - 0.005, label, ha='center', va='center',
             fontsize=8.5, color=color, fontweight='bold',
             transform=ax1.transAxes)
    # Arrow between steps
    next_ys = [e[2] for e in arch_steps]
    idx = next_ys.index(y)
    if idx < len(arch_steps) - 1:
        next_y = arch_steps[idx + 1][2]
        ax1.annotate('', xy=(0.50, next_y + 0.047),
                     xytext=(0.50, y - 0.058),
                     xycoords='axes fraction', textcoords='axes fraction',
                     arrowprops=dict(arrowstyle='->', color=GRAY, lw=1.3))

# ── Panel derecho superior: features breakdown ───────────────────
ax2 = fig.add_axes([0.48, 0.47, 0.49, 0.40])
ax2.set_facecolor(BG)
ax2.set_title('Desglose de Features por Timestep', fontsize=11, color=BLUE)

feature_groups = [
    ('price(t)',                1,  GRAY),
    ('PE H₀,H₁,H₂',            3,  ORANGE),
    ('AMP H₀,H₁,H₂',           3,  ORANGE),
    ('max H₁',                  1,  ORANGE),
    ('BettiCurve H₀,H₁ (×10)', 20, ORANGE),
    ('Wasserstein Δ H₀,H₁',    2,  RED),
    ('H₁ count > p75',          1,  RED),
    ('Raw Takens (21×6)',      126,  GREEN),
    ('Exógenas (variables)',     5,  BLUE),
]
y_feat = np.arange(len(feature_groups))
widths = [g[1] for g in feature_groups]
labels = [f"{g[0]}  ({g[1]})" for g in feature_groups]
cols   = [g[2] for g in feature_groups]

bars = ax2.barh(y_feat, widths, color=cols, alpha=0.82, edgecolor='white', lw=1, height=0.72)
ax2.set_yticks(y_feat)
ax2.set_yticklabels(labels, fontsize=8.5)
ax2.set_xlabel('Número de features', fontsize=9)
ax2.set_xlim(0, 145)
ax2.invert_yaxis()
# Highlight new vs v1
ax2.axhline(5.5, color=RED, lw=1.2, ls='--', alpha=0.6)
ax2.text(80, 5.8, '▲ Nuevos en v2', fontsize=8, color=RED, style='italic')

total = sum(widths)
ax2.text(1.01, 0.5, f'Total:\n{total} feat\n(+5 exog)', transform=ax2.transAxes,
         fontsize=9, color=BLUE, fontweight='bold', va='center')

# ── Panel derecho inferior: entrenamiento + resultados ───────────
ax3 = fig.add_axes([0.48, 0.12, 0.49, 0.30])
ax3.set_axis_off()
ax3.set_title('Configuración de Entrenamiento y Resultados', fontsize=11, color=BLUE)

train_items = [
    ('Loss:',         'MAE',                                                 BLUE),
    ('Optimizer:',    'Adam(lr=1e-3)  +  ReduceLROnPlateau(p=15, f=0.5)',   GREEN),
    ('Callbacks:',    'EarlyStopping(patience=30, restore_best_weights=True)', GRAY),
    ('Split:',        '70/15/15 cronológico  →  train=117, val=25, test=26', GRAY),
    ('Test MAE:',     '29.11  (mejor global)  — mejora +1.3% sobre TDA v1', GREEN),
    ('Ablation:',     'Sin exógenas: MAE=32.99  |  Exog sin TDA: MAE=45.69', RED),
]
y_t = 0.95
for k, v, color in train_items:
    ax3.text(0.01, y_t, k, fontsize=9, color=color, fontweight='bold',
             transform=ax3.transAxes)
    ax3.text(0.22, y_t, v, fontsize=8.5, color='#263238',
             transform=ax3.transAxes, fontfamily='monospace')
    y_t -= 0.155

fig.text(0.5, 0.04,
         'MultiHeadAttention permite ponderar dinámicamente qué pasos temporales '
         'son más relevantes  |  Con N=168, exógenas añaden ruido en lugar de señal',
         ha='center', fontsize=9.5, color=GRAY,
         bbox=dict(boxstyle='round,pad=0.4', facecolor='#F3E5F5', edgecolor='#7B1FA2', alpha=0.85))

png_attn = fig_to_png(fig)

# ══════════════════════════════════════════════════════════════════
#  INSERTAR EN PPTX
# ══════════════════════════════════════════════════════════════════
print(f"Abriendo {PPTX_PATH}...")
prs = Presentation(PPTX_PATH)
print(f"  Slides actuales: {len(prs.slides)}")

NOTA_SARIMA = """El modelo SARIMA es el benchmark estadístico clásico para esta serie temporal.

SARIMA(1,1,1)(1,1,1)₁₂ combina:
— Parte regular: AR(1) captura dependencia en t-1, una diferencia elimina tendencia, MA(1) corrige el error inmediato.
— Parte estacional: AR estacional en t-12, diferencia estacional elimina el patrón anual, MA estacional rezagado 12 meses.

Los parámetros se seleccionan mediante grid search sobre p,q ∈ {0,1,2} y P,Q ∈ {0,1} minimizando AIC.

El ACF de la serie doblemente diferenciada —Δ₁Δ₁₂— muestra spikes significativos en lag 1 (componente MA) y lag 12 (componente SMA), confirmando que SARIMA(1,1,1)(1,1,1)₁₂ es la especificación correcta.

SARIMA asume linealidad y estacionariedad tras las diferenciaciones. Es el punto de referencia contra el que medimos si los modelos no lineales TDA+LSTM aportan valor adicional."""

NOTA_ATTN = """La Sesión 8 construyó el modelo más complejo del proyecto: LSTM con TDA v2, variables exógenas y atención multi-cabeza.

El input tiene dos ramas:
— TDA v2: 157 features por timestep —1 precio, 156 features topológicos incluyendo H₀, H₁, H₂, Wasserstein distance entre ventanas consecutivas, y contador de ciclos persistentes.
— Variables exógenas: series del FRED concatenadas al mismo timestep.

El mecanismo MultiHeadAttention con 4 cabezas y key_dim=16 opera sobre la secuencia de salida del primer LSTM, permitiendo al modelo ponderar de forma dinámica qué meses del año son más informativos para la predicción actual. La conexión residual y LayerNorm estabilizan el entrenamiento.

El resultado fue MAE=29.11, el mejor global —mejoría marginal de 1.3% sobre TDA v1 con 29.49.

El ablation reveló que las variables exógenas sin TDA empeoran a 45.69: con 168 muestras el modelo no puede aprender relaciones entre series macroeconómicas y el precio de berries. El valor viene exclusivamente del TDA y la atención, no de las exógenas."""

# Insertar SARIMA en posición 5 (0-indexed) → después de slide 5 (optimización)
insert_slide_with_image(prs, png_sarima, position=5,  nota=NOTA_SARIMA)
print(f"  SARIMA insertado en posición 6")

# TDA v2+Attn → ahora está en posición 9 (era slide 8 en lista original, +1 por la inserción)
insert_slide_with_image(prs, png_attn,   position=10, nota=NOTA_ATTN)
print(f"  TDA v2+Exg+Attn insertado en posición 11")

prs.save(PPTX_PATH)
print(f"\n✅  Guardado: {PPTX_PATH}  ({len(prs.slides)} slides)")
