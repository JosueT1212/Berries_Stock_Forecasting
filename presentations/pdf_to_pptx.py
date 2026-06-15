"""
Convierte Presentacion-Ejecutiva.pdf → Presentacion-Ejecutiva.pptx
con notas de orador en cada slide.
"""

import io
import fitz                          # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PDF_PATH  = 'Presentacion-Ejecutiva.pdf'
PPTX_PATH = 'Presentacion-Ejecutiva.pptx'

# ── Speaker notes (one string per slide, same order as PDF) ──────────────────
NOTAS = [
    # Slide 1 — Portada
    """Buenos días / buenas tardes. En esta presentación vamos a revisar el análisis completo del Índice de Precios al Productor de berries en Estados Unidos, serie WPUSI01102B del FRED, usando 215 observaciones mensuales desde junio 2008 hasta abril 2026.

El objetivo fue explorar si las técnicas de Análisis Topológico de Datos —TDA— pueden mejorar la predicción de esta serie respecto a modelos estadísticos y redes neuronales LSTM tradicionales.

El trabajo se desarrolló en 9 sesiones: las primeras cuatro de análisis exploratorio y optimización, y las últimas cinco de modelado predictivo progresivo.""",

    # Slide 2 — Los Datos + Subperiodos
    """La serie tiene 215 observaciones mensuales con un rango de precios de aproximadamente 80 a 410 puntos base 1982=100.

Identificamos tres subperiodos con características estadísticas distintas:
— 2008–2012: crisis financiera global, alta volatilidad, precio oscila fuertemente.
— 2013–2019: período de relativa estabilidad, la serie es más predecible.
— 2020–2026: pandemia e inflación post-COVID, tendencia alcista pronunciada.

Esta segmentación fue fundamental para el análisis topológico posterior, ya que la estructura de los diagramas de persistencia cambia significativamente entre subperiodos.""",

    # Slide 3 — ¿Qué es TDA?
    """El Análisis Topológico de Datos se basa en el Teorema de Takens: dada una serie temporal escalar, si reconstruimos el espacio de estados con dimensión D suficiente y retardo τ, obtenemos un punto-nube que preserva la topología del atractor del sistema dinámico original.

Con D=6 y τ=3 meses, cada ventana deslizante de W=36 observaciones genera 21 vectores de 6 dimensiones. La homología persistente de Vietoris-Rips detecta loops H₁ en esa nube: esos loops corresponden a ciclos estacionales en el precio.

Del diagrama de persistencia extraemos features cuantitativos: entropía de persistencia, amplitud, máxima persistencia H₁, curvas de Betti —que cuentan componentes y loops por escala— más el embedding crudo aplanado. Total: 152 features por paso temporal.""",

    # Slide 4 — TDA por Subperiodo
    """La intensidad de los loops H₁ varía significativamente entre subperiodos.

El período 2008–2012 produce los diagramas de persistencia con mayor amplitud H₁. La alta volatilidad de la crisis crea una estructura topológica más pronunciada y distinguible en el espacio de fase.

El período 2013–2019 tiene la menor persistencia H₁, consistente con la estabilidad relativa de los precios en esa época.

El período 2020–2026 muestra valores intermedios: la pandemia recrea volatilidad pero con un patrón topológico diferente al de la crisis de 2008.

Este resultado justifica por qué el TDA captura información estructural que los modelos estadísticos convencionales no aprovechan.""",

    # Slide 5 — Optimización de Parámetros (Sesión 4)
    """En la Sesión 4 realizamos un grid search sistemático sobre los tres parámetros clave del pipeline TDA.

Para el tamaño de ventana W, el óptimo fue 18 meses, que captura aproximadamente 1.5 ciclos estacionales anuales —suficiente para que el embedding forme loops cerrados en el espacio de fase.

Para la dimensión D, el óptimo fue 4, siguiendo la regla práctica D ≤ W/3. Dimensiones mayores reducen los puntos del embedding y pueden subajustar la topología.

Para el retardo τ, el óptimo fue 1 mes para predicción a corto plazo; con τ más largo se pierde autocorrelación local.

Para las sesiones de LSTM adoptamos W=36, D=6, τ=3 que ofrecen mayor estabilidad con la serie temporal completa de 215 observaciones.""",

    # Slide 6 — RF → LSTM+TDA v1 Input/Output
    """La Sesión 6 construyó el primer modelo end-to-end con LSTM y features TDA.

El input tiene forma (117, 12, 152): 117 secuencias de entrenamiento, cada una con 12 pasos temporales —un año de lookback— y 152 features por paso. De esos 152: 1 es el precio corriente, 26 son resúmenes topológicos y 126 son el embedding de Takens aplanado.

La arquitectura es secuencial: LSTM de 64 unidades con return_sequences, LSTM de 32, una capa densa de 16 con activación ReLU, y salida escalar lineal. Se entrenó con MAE loss, Adam optimizer y Early Stopping con paciencia 20.

El resultado: MAE bajó de 38.71 del RandomForest baseline a 29.49 —una mejora del 23.8%— demostrando que el LSTM aprovecha la secuencia temporal de features TDA.""",

    # Slide 7 — LSTM Puro vs LSTM+TDA
    """Esta es la comparación central que responde la pregunta de investigación: ¿agrega valor real el TDA al LSTM?

El LSTM puro, que solo recibe el precio como único feature de entrada —shape (12, 1)— alcanza MAE=34.38.

El LSTM con features TDA v1 llega a 29.49: una mejora del 14.2%. Esto demuestra que los features topológicos aportan información genuina que el LSTM no puede extraer solo de la serie de precios.

El mejor modelo de la Sesión 8, que combina TDA v2 con variables exógenas y mecanismo de atención multi-cabeza, llega a 29.11 —mejoría marginal sobre TDA v1.

El Random Forest baseline, que aplana la ventana completa de features TDA como un vector de 1,824 dimensiones, queda en 38.71, confirmando que la arquitectura secuencial LSTM es esencial para aprovechar los features topológicos.""",

    # Slide 8 — TDA v2 Features
    """En la Sesión 8 enriquecimos el pipeline de features TDA para obtener la versión 2.

Las principales adiciones respecto a v1 fueron:
— Homología H₂: detección de cavidades tridimensionales en el embedding.
— Distancia de Wasserstein entre diagramas de persistencia consecutivos: mide cuánto cambió la topología de una ventana a la siguiente —útil para detectar rupturas estructurales.
— Contador de ciclos H₁ persistentes sobre el percentil 75.

También probamos 7 variables exógenas del FRED y el S&P 500. Sin embargo, con N=168 secuencias de entrenamiento, el modelo no pudo generalizar correctamente esas relaciones. La mayoría de las series exógenas tenían más del 10% de datos faltantes y fueron excluidas automáticamente.""",

    # Slide 9 — Ablation Study
    """El ablation study de la Sesión 8 revela tres hallazgos importantes.

Primero: las variables exógenas perjudican. LSTM más exógenas sin TDA obtiene MAE=45.69, el peor resultado de todos los modelos. Con 168 muestras mensuales no hay suficiente señal para aprender relaciones entre variables macroeconómicas y el precio de berries.

Segundo: TDA v2 sin exógenas tiene MAE=32.99, inferior incluso al TDA v1 original con 29.49. Más features no siempre mejoran con N pequeño; hay riesgo de sobreajuste.

Tercero: combinar TDA v2 con atención multi-cabeza y exógenas da 29.11, ligeramente mejor que TDA v1 solo. Pero la ganancia marginal no justifica la complejidad adicional con este conjunto de datos.

Conclusión del ablation: LSTM+TDA v1 ofrece el mejor balance calidad-complejidad.""",

    # Slide 10 — Sesiones 9 y 10
    """Exploramos dos extensiones avanzadas.

La Sesión 9 implementó EP Loss —Enhanced Peak Loss— que penaliza errores en picos y valles con un factor alpha, más una cabeza de clasificación binaria que predice simultáneamente si el mes es un extremo local. El grid search sobre alpha en {2,3,5,8} y beta en {0.1,0.3,0.5} eligió alpha=5, beta=0.5. Sin embargo, el MAE en test empeoró a 31.33. La razón: con solo 32 extremos locales en 168 muestras —19% de positivos— la señal supervisada de peaks es insuficiente y la penalización asimétrica daña la generalización global.

La Sesión 10 agregó features de mes codificados como seno y coseno del mes objetivo. Esto logró una mejora de 1.63 puntos de MAE. El resultado confirma que el LSTM no estaba infiriendo la estacionalidad anual solo del embedding TDA —la codificación explícita del mes ayuda significativamente.""",

    # Slide 11 — Forecast Visual
    """El gráfico muestra las predicciones del LSTM puro y del LSTM+TDA v1 sobre el conjunto de prueba —las últimas 26 observaciones, aproximadamente dos años.

LSTM+TDA v1 sigue mejor los cambios de dirección de la serie, especialmente en los puntos de inflexión 2024-2025.

Ambos modelos exhiben el efecto de desplazamiento o 'shift': las predicciones tienden a seguir al valor del mes anterior en lugar de anticipar el cambio. Esto es una limitación conocida de los modelos que predicen niveles directamente.

La corrección natural sería entrenar sobre primeras diferencias del precio —el cambio mensual— y reconstruir el nivel con suma acumulada. Esto eliminaría el incentivo del modelo a copiar el último valor observado.""",

    # Slide 12 — Resumen y Conclusiones
    """Para cerrar, el ranking completo de modelos confirma cuatro conclusiones.

Primera: el TDA agrega valor real. La reducción es del 14% en MAE versus LSTM puro y del 23.8% versus el RandomForest baseline.

Segunda: con N=168, la complejidad tiene límites. Más features TDA, variables exógenas y pérdidas asimétricas no mejoran la generalización con este tamaño de muestra.

Tercera: los features de calendario son un quick-win. month_sin y month_cos del mes objetivo reducen MAE en 1.63 puntos con costo computacional mínimo.

Cuarta: el shift implícito es el principal problema no resuelto. El modelo aprende a predecir p(t+1) ≈ p(t), lo que produce predicciones visualmente desplazadas.

El próximo paso natural es entrenar sobre diferencias Δprice con descomposición STL para separar tendencia y estacionalidad antes de modelar. Gracias por su atención.""",
]

# ── Convert PDF pages to PNG images ──────────────────────────────────────────
print("Convirtiendo PDF a imágenes...")
doc   = fitz.open(PDF_PATH)
pages = []
for i, page in enumerate(doc):
    mat  = fitz.Matrix(2.5, 2.5)       # 2.5× = ~240 dpi
    pix  = page.get_pixmap(matrix=mat, alpha=False)
    img_bytes = pix.tobytes('png')
    pages.append(img_bytes)
    print(f"  Página {i+1}/{len(doc)} ✓")
doc.close()

assert len(pages) == len(NOTAS), \
    f"Mismatch: {len(pages)} páginas PDF vs {len(NOTAS)} notas"

# ── Build PPTX ────────────────────────────────────────────────────────────────
print("\nCreando PPTX...")
prs = Presentation()

# 16:9 widescreen
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # totalmente en blanco

for i, (img_bytes, nota) in enumerate(zip(pages, NOTAS)):
    slide = prs.slides.add_slide(blank_layout)

    # Full-bleed image
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(
        img_stream,
        left=0, top=0,
        width=prs.slide_width,
        height=prs.slide_height
    )

    # Speaker notes
    notes_slide = slide.notes_slide
    tf          = notes_slide.notes_text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = nota.strip()
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0x26, 0x32, 0x38)

    print(f"  Slide {i+1:02d} ✓")

prs.save(PPTX_PATH)
print(f"\n✅  Guardado: {PPTX_PATH}")
print(f"    {len(pages)} slides con notas de orador")
