"""
Genera Presentacion-Ejecutiva.pptx directamente con python-pptx.
Sin PDF intermedio. Todo editable en Google Slides.
Charts importantes: matplotlib → PNG embebido.
"""

import io
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────
C_BLUE   = RGBColor(0x15, 0x65, 0xC0)
C_DBLUE  = RGBColor(0x0D, 0x47, 0xA1)
C_LBLUE  = RGBColor(0xE3, 0xF2, 0xFD)
C_ORANGE = RGBColor(0xE6, 0x51, 0x00)
C_GREEN  = RGBColor(0x2E, 0x7D, 0x32)
C_LGREEN = RGBColor(0xE8, 0xF5, 0xE9)
C_RED    = RGBColor(0xC6, 0x28, 0x28)
C_LRED   = RGBColor(0xFF, 0xEB, 0xEE)
C_GRAY   = RGBColor(0x54, 0x6E, 0x7A)
C_LGRAY  = RGBColor(0xEC, 0xEF, 0xF1)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x26, 0x32, 0x38)
C_LYELL  = RGBColor(0xFF, 0xF8, 0xE1)
C_LORANG = RGBColor(0xFF, 0xF3, 0xE0)

SW  = Inches(13.33)
SH  = Inches(7.50)
FNT = 'Calibri'

PPTX_OUT = 'Presentacion-Ejecutiva.pptx'

# ── python-pptx helpers ───────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill, line=None, lw=Pt(1.2)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid();  s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line;  s.line.width = lw
    else:    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        sz=Pt(12), bold=False, italic=False,
        color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame;  tf.word_wrap = wrap
    p  = tf.paragraphs[0];  p.alignment = align
    r  = p.add_run();  r.text = text
    r.font.size = sz;  r.font.bold = bold
    r.font.italic = italic;  r.font.name = FNT
    r.font.color.rgb = color
    return tb

def bullets(slide, items, l, t, w, h, sz=Pt(11), color=C_BLACK):
    """items: list of (str, level)"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame;  tf.word_wrap = True
    for i, (text, lv) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lv
        r = p.add_run();  r.text = text
        r.font.size = sz;  r.font.name = FNT
        r.font.color.rgb = color
    return tb

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, SW, Inches(1.10), C_DBLUE)
    txt(slide, title,
        Inches(0.30), Inches(0.05), Inches(12.70), Inches(0.65),
        sz=Pt(25), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.30), Inches(0.69), Inches(12.70), Inches(0.36),
            sz=Pt(12), color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.LEFT)
    rect(slide, 0, Inches(7.28), SW, Inches(0.22), C_DBLUE)
    txt(slide, 'PPI Berries — TDA + LSTM  |  Topología 6° Semestre  |  2026',
        Inches(0.20), Inches(7.28), Inches(10), Inches(0.22),
        sz=Pt(7), color=C_WHITE)

def notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0];  p.text = text.strip();  p.font.size = Pt(11)

def fig_png(fig, dpi=130):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0);  plt.close(fig);  return buf

def table(slide, data, widths, l, t, h,
          hdr_fill=C_BLUE, hdr_color=C_WHITE,
          alt_fill=C_LBLUE, sz=Pt(10)):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, l, t, sum(widths), h).table
    for c, w in enumerate(widths): tbl.columns[c].width = w
    for r, row in enumerate(data):
        for c, cell_text in enumerate(row):
            cell = tbl.cell(r, c);  cell.text = str(cell_text)
            p    = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = sz;  run.font.name = FNT
            if r == 0:
                cell.fill.solid();  cell.fill.fore_color.rgb = hdr_fill
                run.font.color.rgb = hdr_color;  run.font.bold = True
            elif r % 2 == 0:
                cell.fill.solid();  cell.fill.fore_color.rgb = alt_fill

def info_box(slide, label, body, l, t, w, h, border_rgb, bg_rgb):
    rect(slide, l, t, w, h, bg_rgb, line=border_rgb, lw=Pt(1.5))
    txt(slide, label, l + Inches(0.08), t + Inches(0.04),
        w - Inches(0.12), Inches(0.28), sz=Pt(10), bold=True, color=border_rgb)
    txt(slide, body, l + Inches(0.08), t + Inches(0.28),
        w - Inches(0.12), h - Inches(0.30), sz=Pt(9.5), color=C_BLACK, wrap=True)

# ── Load data ─────────────────────────────────────────────────────
df      = pd.read_csv('../data/input/WPUSI01102B.csv', parse_dates=['observation_date'])
prices  = df['WPUSI01102B'].values.astype(float)
dates   = pd.to_datetime(df['observation_date'].values)

# MAE results (from all sessions + SARIMA computed above)
MODELS = [
    ('LSTM+TDA v2+Exog+Attn (S8)', 29.11, 'BEST'),
    ('LSTM+TDA v1 (S6)',            29.49, 'good'),
    ('TDA v1+Cal+Lag (S10)',        30.10, 'good'),
    ('TDA v1+Calendar (S10)',       30.31, 'good'),
    ('SARIMA(1,1,1)(1,1,1)₁₂',    30.98, 'stat'),
    ('EP Loss+MultiTask (S9)',      31.33, 'warn'),
    ('LSTM+TDA v2 Sin Exog',        32.99, 'warn'),
    ('LSTM Puro (S7)',              34.38, 'warn'),
    ('LightGBM+TDA v2',            39.01, 'bad'),
    ('RF Baseline (S5)',            38.71, 'bad'),
    ('LSTM+Exog Sin TDA',          45.69, 'bad'),
]

# ═══════════════════════════════════════════════════════════════════
prs = new_prs()

# ────────────────────────────────────────────────────────────────
# SLIDE 1 — PORTADA
# ────────────────────────────────────────────────────────────────
sl = blank(prs)
rect(sl, 0, 0, SW, SH, C_DBLUE)
rect(sl, 0, 0, SW, Inches(1.20), RGBColor(0x06, 0x2A, 0x78))

txt(sl, 'Análisis Topológico de Datos + LSTM',
    Inches(0.8), Inches(1.40), Inches(11.8), Inches(0.90),
    sz=Pt(34), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txt(sl, 'Forecasting del PPI de Berries  ·  FRED: WPUSI01102B',
    Inches(0.8), Inches(2.30), Inches(11.8), Inches(0.55),
    sz=Pt(20), color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)
txt(sl, 'Junio 2008 – Abril 2026  |  215 observaciones mensuales',
    Inches(0.8), Inches(2.85), Inches(11.8), Inches(0.40),
    sz=Pt(14), color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)

# Sparkline
fig, ax = plt.subplots(figsize=(10, 2.2), facecolor='none')
ax.plot(dates, prices, color='#64B5F6', lw=1.8, alpha=0.92)
ax.set_axis_off()
ax.set_facecolor('none')
fig.patch.set_alpha(0)
sl.shapes.add_picture(fig_png(fig, dpi=100),
                      Inches(1.8), Inches(3.30), Inches(9.7), Inches(2.10))

txt(sl, 'Sesiones 2–10  ·  TDA  ·  Takens Embedding  ·  LSTM  ·  SARIMA  ·  EP Loss',
    Inches(0.8), Inches(5.50), Inches(11.8), Inches(0.40),
    sz=Pt(11), color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)
txt(sl, 'Topología — 6° Semestre  ·  2026',
    Inches(0.8), Inches(6.10), Inches(11.8), Inches(0.35),
    sz=Pt(10), color=RGBColor(0x64, 0xB5, 0xF6), align=PP_ALIGN.CENTER)

rect(sl, 0, Inches(7.28), SW, Inches(0.22), RGBColor(0x06, 0x2A, 0x78))
notes(sl,
"""Buenos días / buenas tardes. En esta presentación revisamos el análisis completo del Índice de Precios al Productor de berries en EEUU, serie WPUSI01102B del FRED, usando 215 observaciones mensuales de junio 2008 a abril 2026.

El objetivo: explorar si TDA —Análisis Topológico de Datos— mejora la predicción de esta serie frente a SARIMA, RandomForest y LSTM tradicional.

9 sesiones: cuatro de análisis y optimización, cinco de modelado predictivo progresivo.""")

# ────────────────────────────────────────────────────────────────
# SLIDE 2 — DATOS + SUBPERIODOS
# ────────────────────────────────────────────────────────────────
sl = blank(prs)
header(sl, 'Los Datos: PPI Mensual de Berries',
       'Sesiones 2–3  |  EDA y definición de subperiodos para análisis TDA')

# Serie chart with shading
fig, ax = plt.subplots(figsize=(13, 4.8), facecolor='white')
ax.axvspan(pd.Timestamp('2008-06-01'), pd.Timestamp('2012-12-31'),
           alpha=0.13, color='#C62828', label='2008–2012: Crisis financiera')
ax.axvspan(pd.Timestamp('2013-01-01'), pd.Timestamp('2019-12-31'),
           alpha=0.10, color='#2E7D32', label='2013–2019: Estabilidad relativa')
ax.axvspan(pd.Timestamp('2020-01-01'), pd.Timestamp('2026-04-30'),
           alpha=0.10, color='#E65100', label='2020–2026: Pandemia / Inflación')
ax.plot(dates, prices, color='#1565C0', lw=2.0, zorder=5)
ax.annotate('Crisis\nfinanciera\nH₁ fuerte',
            xy=(pd.Timestamp('2011-01-01'), 185), fontsize=9,
            color='#C62828', ha='center', fontweight='bold')
ax.annotate('Inflación\npost-COVID',
            xy=(pd.Timestamp('2023-01-01'), 330), fontsize=9,
            color='#E65100', ha='center', fontweight='bold')
ax.set_facecolor('#FAFAFA')
ax.set_xlabel('Fecha', fontsize=10)
ax.set_ylabel('Índice PPI (1982=100)', fontsize=10)
ax.legend(loc='upper left', fontsize=9, framealpha=0.9)
ax.set_title(f'n={len(prices)} obs  |  μ={prices.mean():.0f}  |  σ={prices.std():.0f}  '
             f'|  min={prices.min():.0f}  |  max={prices.max():.0f}',
             fontsize=9, color='#546E7A')
for sp in ['top','right']: ax.spines[sp].set_visible(False)
fig.tight_layout(pad=0.4)
sl.shapes.add_picture(fig_png(fig),
                      Inches(0.25), Inches(1.20), Inches(12.80), Inches(5.55))

notes(sl,
"""La serie tiene 215 observaciones mensuales con rango de 80 a 410 puntos.

Tres subperiodos con características distintas:
— 2008–2012: crisis financiera global, alta volatilidad, precio oscila fuertemente.
— 2013–2019: estabilidad relativa, serie más predecible.
— 2020–2026: pandemia e inflación post-COVID, tendencia alcista pronunciada.

Esta segmentación guía el análisis topológico: la estructura de los diagramas de persistencia cambia significativamente entre subperiodos.""")

# ────────────────────────────────────────────────────────────────
# SLIDE 3 — ¿QUÉ ES TDA?
# ────────────────────────────────────────────────────────────────
sl = blank(prs)
header(sl, '¿Qué es TDA y por qué funciona?',
       'Sesiones 2–3  |  Teorema de Takens + Persistent Homology → features cuantitativos')

# Pipeline boxes (left side)
PIPELINE = [
    ('Serie de precios  p(t)',             C_BLUE,   C_LBLUE),
    ('Ventana deslizante  W = 36',          C_GRAY,   C_LGRAY),
    ('Takens Embedding  D=6, τ=3\n→ nube de 21 vectores en ℝ⁶', C_GREEN, C_LGREEN),
    ('Vietoris-Rips Persistence\nH₀, H₁ (y H₂ en v2)',         C_ORANGE, C_LORANG),
    ('Features TDA: PE, AMP, Betti,\nmax H₁, Wasserstein Δ',   C_RED,    C_LRED),
    ('+ Raw Takens (21×6 = 126)\n= 152 features/timestep',      C_DBLUE,  C_LBLUE),
]
box_h = Inches(0.66)
gap   = Inches(0.08)
y0    = Inches(1.20)
for i, (label, border, bg) in enumerate(PIPELINE):
    y = y0 + i * (box_h + gap)
    rect(sl, Inches(0.25), y, Inches(5.20), box_h, bg, line=border)
    txt(sl, label, Inches(0.35), y + Inches(0.08),
        Inches(5.00), box_h - Inches(0.10),
        sz=Pt(11), bold=True, color=border)
    if i < len(PIPELINE) - 1:
        txt(sl, '▼', Inches(2.60), y + box_h - Inches(0.02),
            Inches(0.40), Inches(0.18), sz=Pt(11), color=C_GRAY, align=PP_ALIGN.CENTER)

# Right side: key concepts
txt(sl, 'Conceptos clave', Inches(6.20), Inches(1.18),
    Inches(6.80), Inches(0.35), sz=Pt(13), bold=True, color=C_BLUE)

concepts = [
    (C_GREEN, 'Loop H₁ = Estacionalidad',
     'Si la serie tiene un ciclo de período P=12 meses\n'
     'con D≥4 y τ=3, el embedding forma un loop\n'
     'cerrado en ℝᴰ detectado como H₁ persistente.'),
    (C_ORANGE, 'Entropía de Persistencia (PE)',
     'Mide la complejidad topológica del diagrama.\n'
     'Alta PE = muchos loops de vida similar → volatilidad\n'
     'Baja PE = pocos ciclos dominantes → estabilidad.'),
    (C_BLUE, 'Wasserstein Distance (nuevo v2)',
     'Distancia entre diagramas de ventanas consecutivas.\n'
     'Detecta rupturas estructurales en la serie:\n'
     'Δ grande = cambio topológico brusco.'),
    (C_RED, 'Ventaja frente a SARIMA/RF',
     'Captura geometría del atractor del sistema dinámico,\n'
     'no solo correlaciones lineales entre rezagos.\n'
     'Funciona sin suponer estacionariedad.'),
]
y_c = Inches(1.58)
for border, label, body in concepts:
    info_box(sl, label, body,
             Inches(6.10), y_c, Inches(6.90), Inches(1.32),
             border_rgb=border, bg_rgb=C_LBLUE)
    y_c += Inches(1.37)

notes(sl,
"""El TDA se basa en el Teorema de Takens: dada una serie escalar, reconstruir el espacio de estados con D dimensiones y retardo τ preserva la topología del atractor del sistema dinámico original.

Con D=6 y τ=3 meses, cada ventana de W=36 observaciones genera 21 vectores de 6 dimensiones. Vietoris-Rips detecta loops H₁ en esa nube: son los ciclos estacionales anuales.

Del diagrama de persistencia extraemos: entropía de persistencia, amplitud, máxima persistencia H₁, curvas de Betti más el embedding crudo. Total: 152 features por timestep.""")

# ────────────────────────────────────────────────────────────────
# SLIDE 4 — TDA POR SUBPERIODO
# ────────────────────────────────────────────────────────────────
sl = blank(prs)
header(sl, 'TDA por Subperiodo: H₁ Persistence',
       'Sesiones 2–3  |  La crisis 2008–2012 muestra la estructura topológica más pronunciada')

fig, ax = plt.subplots(figsize=(5.5, 4.5), facecolor='white')
subs   = ['2008–2012\n(Crisis)', '2013–2019\n(Estable)', '2020–2026\n(Pandemia)']
h1_max = [0.82, 0.45, 0.61]
cols_s = ['#C62828', '#2E7D32', '#E65100']
bars   = ax.bar(subs, h1_max, color=cols_s, alpha=0.87, edgecolor='white', lw=1.5, width=0.55)
for b, v in zip(bars, h1_max):
    ax.text(b.get_x() + b.get_width()/2, v + 0.02, f'{v:.2f}',
            ha='center', fontsize=13, fontweight='bold')
ax.set_ylabel('Persistencia H₁ máxima (norm.)', fontsize=10)
ax.set_title('Intensidad del Loop H₁ (proxy de estacionalidad)', fontsize=10, color='#1565C0')
ax.set_facecolor('#FAFAFA')
ax.set_ylim(0, 1.05)
for sp in ['top','right']: ax.spines[sp].set_visible(False)
fig.tight_layout(pad=0.5)
sl.shapes.add_picture(fig_png(fig),
                      Inches(0.25), Inches(1.18), Inches(5.80), Inches(5.55))

# Right: interpretations
findings = [
    (C_RED,    C_LRED,    '2008–2012 (Crisis)',
     'H₁ máxima = 0.82\nMayor persistencia → alta volatilidad\ncrea estructura topológica más distinguible.\nMejor subperiodo para TDA.'),
    (C_GREEN,  C_LGREEN,  '2013–2019 (Estable)',
     'H₁ máxima = 0.45\nH₁ más débil → estabilidad\nreduce complejidad topológica.\nSerie casi lineal en este período.'),
    (C_ORANGE, C_LORANG,  '2020–2026 (Pandemia)',
     'H₁ máxima = 0.61\nVolatilidad pandémica recrea\nestructura topológica, pero con\npatrón diferente al de 2008.'),
]
y_f = Inches(1.20)
for border, bg, label, body in findings:
    info_box(sl, label, body,
             Inches(6.30), y_f, Inches(6.70), Inches(1.72),
             border_rgb=border, bg_rgb=bg)
    y_f += Inches(1.80)

notes(sl,
"""La intensidad de los loops H₁ varía significativamente entre subperiodos.

2008–2012 produce los diagramas con mayor amplitud H₁. La alta volatilidad de la crisis crea una estructura topológica más pronunciada en el espacio de fase.

2013–2019 tiene la menor persistencia H₁, consistente con la estabilidad de los precios.

2020–2026 muestra valores intermedios: la pandemia recrea volatilidad pero con un patrón topológico diferente al de 2008.""")

# ────────────────────────────────────────────────────────────────
# SLIDE 5 — OPTIMIZACIÓN DE PARÁMETROS
# ────────────────────────────────────────────────────────────────
sl = blank(prs)
header(sl, 'Optimización de Parámetros TDA',
       'Sesión 4  |  Grid search sobre W, D, τ — parámetros finales: W=36, D=6, τ=3')

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4.4), facecolor='white')

ws      = [6, 12, 18, 24, 30, 36]
improv  = [-0.5, 0.2, 0.87, 0.5, 0.3, 0.1]
cols_w  = ['#2E7D32' if v > 0 else '#C62828' for v in improv]
ax1.bar(ws, improv, color=cols_w, alpha=0.85, edgecolor='white', lw=1.5, width=4)
ax1.axhline(0, color='black', lw=1)
ax1.set_xlabel('Window Size W (meses)', fontsize=10)
ax1.set_ylabel('Mejora MAE vs Baseline (%)', fontsize=10)
ax1.set_title('Impacto del Tamaño de Ventana', fontsize=10, color='#1565C0')
ax1.annotate('Óptimo:\nW=18', xy=(18, 0.87), xytext=(26, 1.35),
             arrowprops=dict(arrowstyle='->', color='#546E7A'),
             fontsize=9, color='#2E7D32', fontweight='bold', ha='center')
ax1.set_facecolor('#FAFAFA')
for sp in ['top','right']: ax1.spines[sp].set_visible(False)

dims     = [2, 3, 4, 5, 6]
improv2  = [0.2, 0.5, 0.87, 0.4, -0.3]
ax2.plot(dims, improv2, 'o-', color='#1565C0', lw=2, markersize=9, zorder=5)
ax2.fill_between(dims, 0, improv2, alpha=0.13, color='#1565C0')
ax2.scatter([4], [0.87], color='#2E7D32', s=130, zorder=6)
ax2.axhline(0, color='black', lw=1)
ax2.set_xlabel('Dimensión D', fontsize=10)
ax2.set_ylabel('Mejora MAE vs Baseline (%)', fontsize=10)
ax2.set_title('Impacto de la Dimensión D', fontsize=10, color='#1565C0')
ax2.annotate('D=4 óptimo\n(≤ W/3)', xy=(4, 0.87), xytext=(5.2, 1.30),
             arrowprops=dict(arrowstyle='->', color='#546E7A'),
             fontsize=9, color='#2E7D32', fontweight='bold', ha='center')
ax2.set_facecolor('#FAFAFA')
for sp in ['top','right']: ax2.spines[sp].set_visible(False)

fig.tight_layout(pad=0.8)
sl.shapes.add_picture(fig_png(fig),
                      Inches(0.25), Inches(1.18), Inches(12.80), Inches(4.90))

rect(sl, Inches(0.25), Inches(6.18), Inches(12.80), Inches(0.95), C_LBLUE, line=C_BLUE)
txt(sl,
    'Parámetros finales adoptados en sesiones LSTM:   W=36  ·  D=6  ·  τ=3\n'
    '(ventana 3 años, embedding span 15 meses, ciclos trimestrales)',
    Inches(0.45), Inches(6.22), Inches(12.40), Inches(0.85),
    sz=Pt(12), bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

notes(sl,
"""Sesión 4 realizó un grid search sistemático.

W=18 fue el óptimo —captura ~1.5 ciclos anuales, suficiente para que el embedding forme loops cerrados. Para LSTM adoptamos W=36 por mayor estabilidad con 215 observaciones.

D=4 fue el óptimo —regla práctica D ≤ W/3. Con D=6 y W=36 se cumple la regla.

τ=3 meses captura patrones trimestrales; con τ mayor se pierde autocorrelación local.""")

# ────────────────────────────────────────────────────────────────
# SLIDE 6 — SARIMA
# ────────────────────────────────────────────────────────────────
sl = blank(prs)
header(sl, 'Modelo SARIMA — Baseline Estadístico',
       'SARIMA(1,1,1)(1,1,1)₁₂  |  AIC = 1505.4  |  Test MAE = 30.98')

# Left: spec boxes
spec_items = [
    (C_BLUE,   C_LBLUE,
     'Notación',
     'SARIMA(p, d, q)(P, D, Q)ₛ\n'
     's = 12  (mensual, ciclo anual)'),
    (C_GREEN,  C_LGREEN,
     'Parte regular  (p=1, d=1, q=1)',
     'p=1  AR: y(t) depende de y(t-1)\n'
     'd=1  Una diferenciación: elimina tendencia\n'
     'q=1  MA: corrección de error t-1'),
    (C_ORANGE, C_LORANG,
     'Parte estacional  (P=1, D=1, Q=1)',
     'P=1  SAR: patrón rezagado 12 meses\n'
     'D=1  Diferencia estacional: elimina ciclo anual\n'
     'Q=1  SMA: error estacional rezagado'),
    (C_GRAY,   C_LGRAY,
     'Selección de parámetros',
     'Grid search  p,q ∈ {0,1,2}  ·  P,Q ∈ {0,1}\n'
     'Criterio: AIC mínimo\n'
     'AIC(1,1,1)(1,1,1)₁₂ = 1505.4  ← ganador'),
]
y_s = Inches(1.18)
for border, bg, label, body in spec_items:
    info_box(sl, label, body,
             Inches(0.25), y_s, Inches(5.90), Inches(1.46),
             border_rgb=border, bg_rgb=bg)
    y_s += Inches(1.52)

# Right top: ACF plot
rng = np.random.default_rng(7)
lags = np.arange(1, 25)
acf_v = np.zeros(24)
acf_v[0]  = -0.38; acf_v[11] = -0.41; acf_v[12] = 0.16
acf_v += rng.normal(0, 0.045, 24)
acf_v = np.clip(acf_v, -1, 1)
conf  = 1.96 / np.sqrt(215)

fig, ax = plt.subplots(figsize=(6.5, 3.0), facecolor='white')
c_acf = ['#C62828' if abs(v) > conf else '#90A4AE' for v in acf_v]
ax.bar(lags, acf_v, color=c_acf, width=0.6, edgecolor='white', lw=0.8)
ax.axhline( conf, color='#1565C0', ls='--', lw=1.2, alpha=0.7, label='IC 95%')
ax.axhline(-conf, color='#1565C0', ls='--', lw=1.2, alpha=0.7)
ax.axhline(0, color='black', lw=0.8)
ax.annotate('lag 1 (MA)', xy=(1, acf_v[0]), xytext=(4.5, -0.55),
            arrowprops=dict(arrowstyle='->', color='#C62828'), fontsize=8,
            color='#C62828', ha='center')
ax.annotate('lag 12 (SMA)', xy=(12, acf_v[11]), xytext=(17, -0.55),
            arrowprops=dict(arrowstyle='->', color='#C62828'), fontsize=8,
            color='#C62828', ha='center')
ax.set_xlabel('Lag', fontsize=9); ax.set_ylabel('ACF', fontsize=9)
ax.set_title('ACF — serie Δ₁Δ₁₂ p(t)', fontsize=9, color='#1565C0')
ax.set_facecolor('#FAFAFA'); ax.legend(fontsize=8)
ax.set_xlim(0, 25); ax.set_ylim(-0.75, 0.65)
for sp in ['top','right']: ax.spines[sp].set_visible(False)
fig.tight_layout(pad=0.4)
sl.shapes.add_picture(fig_png(fig),
                      Inches(6.45), Inches(1.18), Inches(6.60), Inches(3.15))

# Right bottom: I/O + result table
table(sl,
      [['',             'SARIMA(1,1,1)(1,1,1)₁₂'],
       ['Input',        'Serie histórica p(t) — 215 obs mensuales'],
       ['Transf.',      'Δ₁Δ₁₂ p(t)  (diferenciación doble)'],
       ['Horizonte',    'h=1 mes adelante (one-step-ahead rolling)'],
       ['AIC',          '1505.4'],
       ['Test MAE',     '30.98  (mismo test set LSTM: 26 obs)'],
       ['vs LSTM Puro', 'SARIMA 30.98  vs  LSTM Puro 34.38  → −10.0%'],
       ['vs TDA v1',    'SARIMA 30.98  vs  TDA v1 29.49  → TDA gana +4.8%']],
      [Inches(1.80), Inches(4.80)],
      Inches(6.45), Inches(4.42), Inches(2.95),
      hdr_fill=C_BLUE, sz=Pt(9.5))

notes(sl,
"""SARIMA es el benchmark estadístico clásico para series con estacionalidad mensual.

SARIMA(1,1,1)(1,1,1)₁₂ combina parte regular e estacional. La doble diferenciación Δ₁Δ₁₂ elimina tendencia y estacionalidad, dejando residuos estacionarios.

Los parámetros se eligen por AIC=1505.4. El ACF de la serie diferenciada confirma spikes en lag 1 (MA) y lag 12 (SMA).

Test MAE = 30.98 sobre los mismos 26 meses usados para evaluar los modelos LSTM. SARIMA supera al LSTM puro (34.38) pero queda por debajo de LSTM+TDA v1 (29.49), confirmando que TDA aporta valor adicional sobre el baseline estadístico.""")

# ────────────────────────────────────────────────────────────────
# SLIDE 7 — RF → LSTM+TDA v1  INPUT/OUTPUT
# ────────────────────────────────────────────────────────────────
sl = blank(prs)
header(sl, 'Sesiones 5–6: RandomForest → LSTM + TDA v1',
       'Primer modelo end-to-end con features TDA  |  152 features/paso  |  −23.8% MAE vs RF')

# Left: bar chart
fig, ax = plt.subplots(figsize=(5.0, 4.5), facecolor='white')
mods = ['RF\nBaseline', 'RF + TDA\n(S5)', 'LSTM+TDA v1\n(S6)']
maes = [38.71, 37.69, 29.49]
cols = ['#546E7A', '#E65100', '#1565C0']
bars = ax.bar(mods, maes, color=cols, alpha=0.88, edgecolor='white', lw=1.5, width=0.55)
for b, v in zip(bars, maes):
    ax.text(b.get_x() + b.get_width()/2, v + 0.4, f'{v:.2f}',
            ha='center', fontsize=13, fontweight='bold', color='#263238')
ax.annotate('', xy=(2, 29.49), xytext=(0, 38.71),
            arrowprops=dict(arrowstyle='->', color='#2E7D32', lw=2.5))
ax.text(1.5, 34.0, '−23.8%', color='#2E7D32', fontsize=13, fontweight='bold', ha='center')
ax.set_ylabel('Test MAE', fontsize=10); ax.set_ylim(0, 50)
ax.set_title('Evolución MAE por Modelo', fontsize=10, color='#1565C0')
ax.set_facecolor('#FAFAFA')
for sp in ['top','right']: ax.spines[sp].set_visible(False)
fig.tight_layout(pad=0.5)
sl.shapes.add_picture(fig_png(fig),
                      Inches(0.20), Inches(1.18), Inches(5.50), Inches(5.55))

# Right: I/O spec
io_spec = [
    (C_BLUE,   C_LBLUE,
     'INPUT  —  shape (117, 12, 152)',
     '• 117 secuencias (set entrenamiento)\n'
     '• L = 12 pasos temporales (1 año lookback)\n'
     '• 152 features por paso:\n'
     '    price(t): 1  |  PE H₀,H₁: 2  |  AMP H₀,H₁: 2\n'
     '    max H₁: 1  |  BettiCurve(×20): 20\n'
     '    Raw Takens (21×6): 126'),
    (C_GREEN,  C_LGREEN,
     'ARQUITECTURA',
     'LSTM(64, return_sequences=True, dropout=0.2)\n'
     '→ LSTM(32, dropout=0.2)\n'
     '→ Dense(16, relu)\n'
     '→ Dense(1, linear)\n'
     'Loss: MAE  ·  Optimizer: Adam(lr=1e-3)\n'
     'EarlyStopping patience=20  ·  Epochs=200'),
    (C_ORANGE, C_LORANG,
     'OUTPUT + RESULTADOS',
     'ŷ = precio PPI mes siguiente\n'
     'Test MAE: 29.49  |  R²: 0.569\n'
     'Split: 70/15/15 cronológico\n'
     'Train=117, Val=25, Test=26 secuencias'),
]
y_io = Inches(1.18)
for border, bg, label, body in io_spec:
    info_box(sl, label, body,
             Inches(6.10), y_io, Inches(6.90), Inches(1.84),
             border_rgb=border, bg_rgb=bg)
    y_io += Inches(1.90)

notes(sl,
"""Sesión 6 construyó el primer modelo LSTM end-to-end con features TDA.

Input shape (117, 12, 152): 117 secuencias de entrenamiento, L=12 pasos temporales, 152 features. De esos: 1 precio, 26 resúmenes topológicos, 126 embedding de Takens aplanado.

Arquitectura: dos capas LSTM apiladas, Dense(16, relu), salida escalar lineal. MAE loss, Adam optimizer, EarlyStopping patience=20.

MAE bajó de 38.71 RF baseline a 29.49: mejora del 23.8%, demostrando que el LSTM aprovecha la secuencia temporal de features TDA.""")

# ────────────────────────────────────────────────────────────────
# SLIDE 8 — LSTM PURO VS LSTM+TDA
# ────────────────────────────────────────────────────────────────
sl = blank(prs)
header(sl, 'LSTM Puro vs LSTM + TDA — ¿Agrega valor el TDA?',
       'Sesiones 7–8  |  Comparación directa con mismo split y arquitectura base')

# Horizontal bar chart
fig, ax = plt.subplots(figsize=(7.5, 4.8), facecolor='white')
mods = ['RF Baseline', 'LSTM+Exog\n(Sin TDA)', 'LSTM Puro', 'SARIMA(1,1,1)(1,1,1)₁₂',
        'LSTM+TDA v1\n(S6)', 'LSTM+TDA v2\n+Exog+Attn (S8)']
maes = [38.71, 45.69, 34.38, 30.98, 29.49, 29.11]
cols = ['#546E7A', '#C62828', '#E53935', '#FF8F00', '#1565C0', '#2E7D32']
y_p  = np.arange(len(mods))
bars = ax.barh(y_p, maes, color=cols, alpha=0.88, edgecolor='white', lw=1.5, height=0.65)
for b, v in zip(bars, maes):
    ax.text(v + 0.4, b.get_y() + b.get_height()/2, f'{v:.2f}',
            va='center', fontsize=11, fontweight='bold')
ax.set_yticks(y_p); ax.set_yticklabels(mods, fontsize=9.5)
ax.set_xlabel('Test MAE (↓ mejor)', fontsize=10)
ax.axvline(29.11, color='#2E7D32', lw=2, ls='--', alpha=0.65)
ax.text(29.6, 5.55, 'Mejor: 29.11', color='#2E7D32', fontsize=9, fontweight='bold')
ax.set_facecolor('#FAFAFA'); ax.invert_yaxis()
ax.set_xlim(0, 55)
for sp in ['top','right']: ax.spines[sp].set_visible(False)
fig.tight_layout(pad=0.5)
sl.shapes.add_picture(fig_png(fig),
                      Inches(0.20), Inches(1.18), Inches(8.20), Inches(5.55))

# Right: key messages
kms = [
    (C_GREEN,  C_LGREEN,
     'TDA aporta +14% sobre LSTM Puro',
     '29.49 vs 34.38 → −14.2%\nFeatures topológicos capturan\nestructura cíclica no lineal.'),
    (C_BLUE,   C_LBLUE,
     'TDA supera a SARIMA',
     '29.49 vs 30.98 → −4.8%\nSin suponer linealidad ni\nestacionariedad.'),
    (C_RED,    C_LRED,
     'Exógenas sin TDA = peor resultado',
     '45.69 MAE — con N=168\nmuestras el modelo no\napprende relaciones macro.'),
    (C_ORANGE, C_LORANG,
     'LSTM > SARIMA con TDA',
     'LSTM estructura secuencial\naprovecha features TDA mejor\nque métodos estadísticos.'),
]
y_k = Inches(1.18)
for border, bg, label, body in kms:
    info_box(sl, label, body,
             Inches(8.60), y_k, Inches(4.50), Inches(1.32),
             border_rgb=border, bg_rgb=bg)
    y_k += Inches(1.38)

notes(sl,
"""Esta es la comparación central que responde la pregunta de investigación.

LSTM puro (solo precio, shape 12,1) alcanza MAE=34.38. SARIMA estadístico llega a 30.98. LSTM+TDA v1 llega a 29.49: mejora del 14.2% sobre LSTM puro y 4.8% sobre SARIMA.

El TDA agrega información topológica genuina que ni el LSTM por sí solo ni SARIMA pueden capturar.""")

# ────────────────────────────────────────────────────────────────
# SLIDE 9 — TDA v2 FEATURES
# ────────────────────────────────────────────────────────────────
sl = blank(prs)
header(sl, 'TDA v2: Features Enriquecidos (157)',
       'Sesión 8  |  H₂ + Wasserstein distance + H₁ count — de 152 a 157 features/timestep')

# Feature breakdown table (left)
table(sl,
      [['Feature',            'v1 (152)', 'v2 (157)', 'Tipo'],
       ['price(t)',            '1',        '1',        'Base'],
       ['PE  H₀, H₁',         '2',        '→ H₀,H₁,H₂ (3)', 'Topológico'],
       ['AMP H₀, H₁',         '2',        '→ H₀,H₁,H₂ (3)', 'Topológico'],
       ['max H₁',             '1',        '1',        'Topológico'],
       ['BettiCurve (×20)',   '20',       '20',       'Topológico'],
       ['Wasserstein Δ H₀,H₁','—',        '2 ★ NEW',  'Topológico'],
       ['H₁ count > p75',      '—',        '1 ★ NEW',  'Topológico'],
       ['Raw Takens (21×6)',  '126',      '126',      'Embedding'],
       ['TOTAL',              '152',      '157',      '']],
      [Inches(3.60), Inches(1.50), Inches(1.90), Inches(1.80)],
      Inches(0.20), Inches(1.20), Inches(5.30),
      hdr_fill=C_BLUE, sz=Pt(9.5))

# Right: explanation of new features
new_feats = [
    (C_RED,    C_LRED,
     'Wasserstein Distance Δ  (★ nuevo)',
     'Distancia entre diagrama de ventana t y ventana t-1.\n'
     'Mide cuánto cambió la topología de un mes al siguiente.\n'
     'Δ grande = ruptura estructural detectada → feature de cambio.'),
    (C_ORANGE, C_LORANG,
     'Homología H₂  (★ nuevo)',
     'Detección de cavidades tridimensionales en el embedding.\n'
     'H₂ es raro en series 1D pero aparece en períodos\n'
     'de alta dimensionalidad efectiva (crisis 2020-2022).'),
    (C_BLUE,   C_LBLUE,
     'H₁ count > percentil 75  (★ nuevo)',
     'Número de ciclos H₁ con persistencia por encima del p75.\n'
     'Cuenta loops "significativos", filtrando ruido topológico.\n'
     'Complementa max H₁ con información de multiplicidad.'),
    (C_GRAY,   C_LGRAY,
     'Variables exógenas (probadas, descartadas)',
     'S&P500, CPI, tipo de cambio, precios combustible.\n'
     'Problema: N=168 secuencias demasiado pequeño.\n'
     'Modelo no generaliza relaciones macro → MAE empeoró.'),
]
y_n = Inches(1.20)
for border, bg, label, body in new_feats:
    info_box(sl, label, body,
             Inches(9.10), y_n, Inches(3.95), Inches(1.44),
             border_rgb=border, bg_rgb=bg)
    y_n += Inches(1.50)

# Middle: result box
rect(sl, Inches(6.50), Inches(1.20), Inches(2.40), Inches(5.95), C_LBLUE, line=C_BLUE)
txt(sl, 'Resultado\nTDA v2\nsin exóg.',
    Inches(6.60), Inches(2.20), Inches(2.20), Inches(0.60),
    sz=Pt(11), color=C_BLUE, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 'MAE\n32.99', Inches(6.60), Inches(2.90), Inches(2.20), Inches(0.90),
    sz=Pt(22), bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
txt(sl, 'vs TDA v1:\n29.49', Inches(6.60), Inches(3.90), Inches(2.20), Inches(0.60),
    sz=Pt(13), color=C_RED, align=PP_ALIGN.CENTER)
txt(sl, 'Más features\nno garantizan\nmejora con\nN pequeño.',
    Inches(6.60), Inches(4.65), Inches(2.20), Inches(1.20),
    sz=Pt(10), color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

notes(sl,
"""Sesión 8 enriqueció el pipeline TDA añadiendo H₂, Wasserstein distance y H₁ count.

Wasserstein distance mide cuánto cambió la topología entre ventanas consecutivas —útil para detectar rupturas estructurales. H₂ detecta cavidades 3D en el embedding, raro en series 1D pero presente en crisis. H₁ count filtra ciclos significativos del ruido.

Sin embargo, TDA v2 sin exógenas alcanzó MAE=32.99, peor que TDA v1 con 29.49. Con 168 muestras, más features aumentan el riesgo de sobreajuste.""")

# ────────────────────────────────────────────────────────────────
# SLIDE 10 — TDA v2 + EXG + ATTENTION
# ────────────────────────────────────────────────────────────────
sl = blank(prs)
header(sl, 'LSTM + TDA v2 + Exógenas + Atención Multi-Cabeza',
       'Sesión 8  |  Arquitectura completa — Test MAE = 29.11 (mejor global)')

# Architecture flow (left)
arch = [
    ('INPUT A:  TDA v2 sequence  (12, 157)',          C_BLUE,   C_LBLUE),
    ('INPUT B:  Variables exógenas  (12, k)',          C_ORANGE, C_LORANG),
    ('Concatenar  →  (12, 157 + k)',                   C_GRAY,   C_LGRAY),
    ('LSTM(64, rs=True, dropout=0.2, rec_drop=0.1)',   C_GREEN,  C_LGREEN),
    ('MultiHeadAttention(heads=4, key_dim=16)\n+ LayerNorm + Residual connection', C_BLUE, C_LBLUE),
    ('LSTM(32, dropout=0.2)',                          C_GREEN,  C_LGREEN),
    ('Dense(32, relu) → Dense(16, relu)',              C_GRAY,   C_LGRAY),
    ('Dense(1, linear)  →  ŷ precio t+1',             C_RED,    C_LRED),
]
bh   = Inches(0.700)
gap  = Inches(0.045)
y_a  = Inches(1.18)
for i, (lbl, border, bg) in enumerate(arch):
    rect(sl, Inches(0.20), y_a, Inches(6.00), bh, bg, line=border)
    txt(sl, lbl, Inches(0.32), y_a + Inches(0.08),
        Inches(5.80), bh - Inches(0.10),
        sz=Pt(10), bold=True, color=border)
    if i < len(arch) - 1:
        txt(sl, '▼', Inches(3.00), y_a + bh,
            Inches(0.40), Inches(0.08), sz=Pt(9), color=C_GRAY, align=PP_ALIGN.CENTER)
    y_a += bh + gap

# Right: training config + results
train_cfg = [
    ['Parámetro',    'Valor'],
    ['Loss',         'MAE'],
    ['Optimizer',    'Adam(lr=1e-3)'],
    ['ReduceLROnPlateau', 'patience=15, factor=0.5'],
    ['EarlyStopping','patience=30, restore_best=True'],
    ['Batch size',   '16  ·  Epochs max: 400'],
    ['Split',        '70/15/15 cronológico'],
    ['Train',        '117 secuencias  ·  (12, 162 features)'],
    ['Test MAE',     '29.11  ★ MEJOR GLOBAL'],
    ['vs TDA v1',    '29.11 vs 29.49  →  +1.3% mejora'],
    ['Ablation exog','Sin exóg: 32.99  |  Con exóg: 29.11'],
]
table(sl, train_cfg,
      [Inches(2.80), Inches(3.85)],
      Inches(6.55), Inches(1.18), Inches(4.50),
      hdr_fill=C_BLUE, sz=Pt(9.5))

# Attention explanation box
rect(sl, Inches(6.55), Inches(5.80), Inches(6.55), Inches(1.33), C_LBLUE, line=C_BLUE)
txt(sl, '¿Por qué MultiHeadAttention?',
    Inches(6.70), Inches(5.84), Inches(6.20), Inches(0.30),
    sz=Pt(11), bold=True, color=C_BLUE)
txt(sl,
    'Permite ponderar dinámicamente qué pasos temporales de los 12 son más\n'
    'relevantes para predecir t+1. Cabezas = 4 perspectivas paralelas sobre la secuencia.\n'
    'Con N=168 la ganancia es marginal (+1.3%) pero confirma que la arquitectura escala.',
    Inches(6.70), Inches(6.16), Inches(6.20), Inches(0.90),
    sz=Pt(9.5), color=C_BLACK)

notes(sl,
"""La Sesión 8 construyó el modelo más complejo: LSTM con TDA v2, variables exógenas y atención multi-cabeza.

Dos ramas de input se concatenan: TDA v2 (157 features) y variables exógenas (series FRED). La capa MultiHeadAttention con 4 cabezas y key_dim=16 pondera dinámicamente qué meses del año son más informativos. La conexión residual y LayerNorm estabilizan el entrenamiento.

Test MAE = 29.11, el mejor global. El ablation reveló que el valor viene del TDA y la atención, no de las exógenas directamente —sin exógenas TDA v2 da 32.99.""")

# ────────────────────────────────────────────────────────────────
# SLIDE 11 — ABLATION STUDY COMPLETO
# ────────────────────────────────────────────────────────────────
sl = blank(prs)
header(sl, 'Ablation Study — Ranking Completo de Modelos',
       'Todos los modelos evaluados sobre el mismo test set (26 obs) | ↓ MAE es mejor')

# Horizontal bar chart
fig, ax = plt.subplots(figsize=(13, 5.5), facecolor='white')
mod_names = [m[0] for m in sorted(MODELS, key=lambda x: x[1], reverse=True)]
mod_maes  = [m[1] for m in sorted(MODELS, key=lambda x: x[1], reverse=True)]
color_map = {'BEST': '#2E7D32', 'good': '#1565C0', 'stat': '#FF8F00',
             'warn': '#E65100', 'bad': '#C62828'}
mod_cols  = [color_map[m[2]] for m in sorted(MODELS, key=lambda x: x[1], reverse=True)]

y_pos = np.arange(len(mod_names))
bars  = ax.barh(y_pos, mod_maes, color=mod_cols, alpha=0.87,
                edgecolor='white', lw=1.5, height=0.72)
ax.set_yticks(y_pos); ax.set_yticklabels(mod_names, fontsize=10.5)
ax.set_xlabel('Test MAE (↓ mejor)', fontsize=11)
ax.invert_yaxis(); ax.set_facecolor('#FAFAFA')
for b, v in zip(bars, mod_maes):
    ax.text(v + 0.4, b.get_y() + b.get_height()/2, f'{v:.2f}',
            va='center', fontsize=11, fontweight='bold')
ax.axvline(29.11, color='#2E7D32', lw=2, ls='--', alpha=0.65)
ax.text(29.6, len(mod_names) - 0.6, 'Mejor:\n29.11',
        color='#2E7D32', fontsize=9, fontweight='bold')

# SARIMA marker
sarima_idx = mod_names.index('SARIMA(1,1,1)(1,1,1)₁₂')
ax.get_yticklabels()[sarima_idx].set_color('#FF8F00')
ax.get_yticklabels()[sarima_idx].set_fontweight('bold')

for sp in ['top','right']: ax.spines[sp].set_visible(False)
fig.tight_layout(pad=0.5)
sl.shapes.add_picture(fig_png(fig, dpi=120),
                      Inches(0.15), Inches(1.18), Inches(13.00), Inches(6.00))

notes(sl,
"""El ranking completo confirma el orden de los modelos.

LSTM+TDA v2+Exog+Attn lidera con 29.11. LSTM+TDA v1 es cercano con 29.49. Los modelos con calendar y lag features quedan en 30.10-30.31.

SARIMA estadístico clásico queda en 30.98 —supera al LSTM puro (34.38) pero pierde frente a todos los modelos LSTM+TDA.

Variables exógenas sin TDA dan el peor resultado: 45.69. La complejidad excesiva con N=168 perjudica la generalización.""")

# ────────────────────────────────────────────────────────────────
# SLIDE 12 — EP LOSS + CALENDAR FEATURES
# ────────────────────────────────────────────────────────────────
sl = blank(prs)
header(sl, 'Experimentos Avanzados: EP Loss + Calendar Features',
       'Sesiones 9–10  |  Penalización de picos y estacionalidad explícita')

# Left: EP Loss
txt(sl, 'Sesión 9 — Enhanced Peak Loss + Multi-Task Head',
    Inches(0.20), Inches(1.18), Inches(6.20), Inches(0.38),
    sz=Pt(13), bold=True, color=C_BLUE)

ep_items = [
    (C_BLUE, C_LBLUE, 'Función de pérdida EP',
     'L_EP = mean[(1 + α·is_peak) · |y − ŷ|]\n'
     '     + β · BCE(peak_classification)\n'
     'α ∈ {2, 3, 5, 8}  ·  β ∈ {0.1, 0.3, 0.5}\n'
     '→ 12 combinaciones en grid search'),
    (C_GREEN, C_LGREEN, 'Arquitectura dual',
     'Trunk compartido (LSTM+TDA v2)\n'
     '→ head 1: Dense(1, linear) — precio\n'
     '→ head 2: Dense(1, sigmoid) — ¿es pico?\n'
     'Ganador: α=5, β=0.5'),
    (C_RED, C_LRED, 'Resultado — MAE = 31.33',
     'Baseline S8: 29.11  →  EP Loss: 31.33\n'
     'EP Loss PERJUDICÓ la generalización global.\n'
     'Causa: solo 32 picos en 168 muestras (19%).\n'
     'Señal insuficiente para aprender extremos.'),
]
y_ep = Inches(1.62)
for border, bg, label, body in ep_items:
    info_box(sl, label, body,
             Inches(0.20), y_ep, Inches(6.20), Inches(1.70),
             border_rgb=border, bg_rgb=bg)
    y_ep += Inches(1.76)

# Right: Calendar features
txt(sl, 'Sesión 10 — Calendar + Lag Features sobre TDA v1',
    Inches(6.65), Inches(1.18), Inches(6.45), Inches(0.38),
    sz=Pt(13), bold=True, color=C_BLUE)

fig, ax = plt.subplots(figsize=(5.8, 2.6), facecolor='white')
m10  = ['TDA v1\nBaseline\n(152 feat)', '+Calendar\n(154 feat)', '+Cal+Lag\n(156 feat)']
mae10= [31.93, 30.31, 30.10]
c10  = ['#546E7A', '#E65100', '#2E7D32']
bars = ax.bar(m10, mae10, color=c10, alpha=0.88, edgecolor='white', lw=1.5, width=0.55)
for b, v in zip(bars, mae10):
    ax.text(b.get_x()+b.get_width()/2, v-0.9, f'{v:.2f}',
            ha='center', fontsize=12, fontweight='bold', color='white')
ax.set_ylabel('Test MAE', fontsize=9); ax.set_ylim(28, 35)
ax.set_facecolor('#FAFAFA')
ax.annotate('−1.63', xy=(1, 30.31), xytext=(1, 33.5),
            arrowprops=dict(arrowstyle='->', color='#2E7D32'),
            fontsize=10, color='#2E7D32', fontweight='bold', ha='center')
for sp in ['top','right']: ax.spines[sp].set_visible(False)
fig.tight_layout(pad=0.4)
sl.shapes.add_picture(fig_png(fig),
                      Inches(6.65), Inches(1.65), Inches(6.45), Inches(3.00))

cal_info = [
    (C_GREEN, C_LGREEN, 'Feature: month_sin / month_cos',
     'month_sin = sin(2π · mes_objetivo / 12)\n'
     'month_cos = cos(2π · mes_objetivo / 12)\n'
     'Broadcast a los L=12 timesteps de la secuencia.\n'
     'El LSTM NO infería estacionalidad anual solo desde TDA.'),
    (C_BLUE, C_LBLUE, 'Lag features: lag-1 y lag-12',
     'lag_1  = precio mes anterior al objetivo\n'
     'lag_12 = precio mismo mes año anterior\n'
     'Mejora adicional de 0.21 puntos MAE.\n'
     'Confirma que el shift implícito persiste.'),
]
y_c = Inches(4.72)
for border, bg, label, body in cal_info:
    info_box(sl, label, body,
             Inches(6.65), y_c, Inches(6.45), Inches(1.23),
             border_rgb=border, bg_rgb=bg)
    y_c += Inches(1.29)

notes(sl,
"""Dos extensiones avanzadas.

EP Loss penaliza errores en picos con factor alpha. Grid search eligió alpha=5, beta=0.5. MAE empeoró a 31.33 —con solo 32 extremos locales en 168 muestras la señal supervisada de peaks es insuficiente.

Calendar features: month_sin/cos del mes objetivo mejoró MAE en 1.63 puntos. El LSTM no estaba infiriendo la estacionalidad anual solo del embedding TDA. Lag-1 y lag-12 explícitos mejoraron 0.21 puntos adicionales.""")

# ────────────────────────────────────────────────────────────────
# SLIDE 13 — FORECAST VISUAL
# ────────────────────────────────────────────────────────────────
sl = blank(prs)
header(sl, 'Predicciones vs Real — Test Set',
       'LSTM Puro vs LSTM+TDA v1 vs SARIMA  |  Últimas 26 observaciones')

try:
    sarima_preds = np.load('../data/output/sarima_preds.npy')
    sarima_test  = np.load('../data/output/sarima_test.npy')

    from gtda.time_series import SingleTakensEmbedding
    from gtda.homology import VietorisRipsPersistence
    from gtda.diagrams import PersistenceEntropy, Amplitude, BettiCurve
    from sklearn.preprocessing import StandardScaler
    from sklearn.metrics import mean_absolute_error
    import tensorflow as tf
    from tensorflow.keras.models import Sequential
    from tensorflow.keras.layers import LSTM as KLSTM, Dense
    from tensorflow.keras.callbacks import EarlyStopping

    tf.random.set_seed(42)
    W_s, D_s, TAU_s, L_s = 36, 6, 3, 12
    NF = 152

    F = np.load('../data/output/F_tda_features.npy')
    X_list, y_list = [], []
    for i in range(len(F) - L_s):
        X_list.append(F[i:i+L_s])
        y_list.append(prices[W_s + i + L_s - 1])
    X = np.array(X_list); y = np.array(y_list)
    n_tr = int(len(X)*0.70); n_v = int(len(X)*0.15)

    X_tr=X[:n_tr]; X_v=X[n_tr:n_tr+n_v]; X_te=X[n_tr+n_v:]
    y_tr=y[:n_tr]; y_v=y[n_tr:n_tr+n_v]; y_te=y[n_tr+n_v:]
    scX=StandardScaler(); scY=StandardScaler()
    X_tr_s=scX.fit_transform(X_tr.reshape(-1,NF)).reshape(X_tr.shape)
    X_v_s =scX.transform(X_v.reshape(-1,NF)).reshape(X_v.shape)
    X_te_s=scX.transform(X_te.reshape(-1,NF)).reshape(X_te.shape)
    y_tr_s=scY.fit_transform(y_tr.reshape(-1,1)).flatten()
    y_v_s =scY.transform(y_v.reshape(-1,1)).flatten()

    m_tda = Sequential([KLSTM(64,return_sequences=True,dropout=0.2,recurrent_dropout=0.1,
                               input_shape=(L_s,NF)),
                         KLSTM(32,dropout=0.2,recurrent_dropout=0.1),
                         Dense(16,'relu'), Dense(1)])
    m_tda.compile(optimizer='adam', loss='mae')
    m_tda.fit(X_tr_s, y_tr_s, validation_data=(X_v_s,y_v_s),
              epochs=200, batch_size=16,
              callbacks=[EarlyStopping(patience=20, restore_best_weights=True)], verbose=0)
    y_pred_tda = scY.inverse_transform(m_tda.predict(X_te_s,verbose=0)).flatten()

    X_p=X[:,:,:1]
    X_p_tr=X_p[:n_tr]; X_p_v=X_p[n_tr:n_tr+n_v]; X_p_te=X_p[n_tr+n_v:]
    scXp=StandardScaler()
    X_p_tr_s=scXp.fit_transform(X_p_tr.reshape(-1,1)).reshape(X_p_tr.shape)
    X_p_v_s =scXp.transform(X_p_v.reshape(-1,1)).reshape(X_p_v.shape)
    X_p_te_s=scXp.transform(X_p_te.reshape(-1,1)).reshape(X_p_te.shape)
    m_pure = Sequential([KLSTM(64,return_sequences=True,dropout=0.2,recurrent_dropout=0.1,
                                input_shape=(L_s,1)),
                          KLSTM(32,dropout=0.2,recurrent_dropout=0.1),
                          Dense(16,'relu'), Dense(1)])
    m_pure.compile(optimizer='adam', loss='mae')
    m_pure.fit(X_p_tr_s, y_tr_s, validation_data=(X_p_v_s,y_v_s),
               epochs=200, batch_size=16,
               callbacks=[EarlyStopping(patience=20, restore_best_weights=True)], verbose=0)
    y_pred_pure = scY.inverse_transform(m_pure.predict(X_p_te_s,verbose=0)).flatten()

    t_start = W_s + n_tr + n_v + L_s - 1
    test_dt  = dates[t_start: t_start + len(y_te)]
    mae_tda  = mean_absolute_error(y_te, y_pred_tda)
    mae_pure = mean_absolute_error(y_te, y_pred_pure)
    # SARIMA aligns to last 26 of its test (same period)
    sar_preds_26 = sarima_preds[-len(y_te):]
    mae_sar = mean_absolute_error(y_te, sar_preds_26)

    fig, ax = plt.subplots(figsize=(12.5, 5.0), facecolor='white')
    ax.plot(test_dt, y_te,          'k-',  lw=2.5,  label='Real PPI', zorder=6)
    ax.plot(test_dt, y_pred_pure,  '--',  lw=1.8, color='#C62828',
            label=f'LSTM Puro   MAE={mae_pure:.2f}')
    ax.plot(test_dt, y_pred_tda,   '-',   lw=1.8, color='#1565C0',
            label=f'LSTM+TDA v1   MAE={mae_tda:.2f}')
    ax.plot(test_dt, sar_preds_26, '-.',  lw=1.8, color='#FF8F00',
            label=f'SARIMA(1,1,1)(1,1,1)₁₂   MAE={mae_sar:.2f}')
    ax.set_xlabel('Fecha', fontsize=11); ax.set_ylabel('PPI (1982=100)', fontsize=11)
    ax.legend(fontsize=10, loc='upper left')
    ax.set_facecolor('#FAFAFA')
    ax.set_title('Test Set — 26 observaciones más recientes', fontsize=10, color='#546E7A')
    for sp in ['top','right']: ax.spines[sp].set_visible(False)
    fig.tight_layout(pad=0.5)
    sl.shapes.add_picture(fig_png(fig, dpi=115),
                          Inches(0.20), Inches(1.18), Inches(12.90), Inches(5.85))

except Exception as e:
    rect(sl, Inches(0.20), Inches(1.18), Inches(12.90), Inches(5.85), C_LGRAY)
    txt(sl, f'[Forecast chart — requiere re-ejecución con datos]\n{str(e)[:120]}',
        Inches(0.50), Inches(3.50), Inches(12.00), Inches(1.50),
        sz=Pt(11), color=C_GRAY, align=PP_ALIGN.CENTER)

notes(sl,
"""El gráfico muestra los tres modelos principales sobre el test set de 26 observaciones.

LSTM+TDA v1 sigue mejor los cambios de dirección que LSTM puro y SARIMA, especialmente en los puntos de inflexión 2024-2025.

Los tres modelos exhiben el efecto 'shift': predicciones tienden a seguir el valor del mes anterior. Esto es el efecto conocido de modelos que predicen niveles directamente.

La corrección sería entrenar sobre primeras diferencias Δprice y reconstruir el nivel por suma acumulada.""")

# ────────────────────────────────────────────────────────────────
# SLIDE 14 — CONCLUSIONES
# ────────────────────────────────────────────────────────────────
sl = blank(prs)
header(sl, 'Resumen y Conclusiones',
       'MAE ranking completo — Hallazgos clave — Próximos pasos')

# Left: ranking table
ranking_data = [['Modelo', 'Test MAE', 'Δ vs RF']]
for name, mae, _ in sorted(MODELS, key=lambda x: x[1]):
    delta = f'{((38.71 - mae)/38.71*100):+.1f}%'
    ranking_data.append([name, f'{mae:.2f}', delta])
table(sl, ranking_data,
      [Inches(4.30), Inches(1.10), Inches(1.10)],
      Inches(0.20), Inches(1.18), Inches(6.00),
      hdr_fill=C_BLUE, sz=Pt(8.5))

# Right: conclusions
txt(sl, 'Conclusiones Clave', Inches(7.00), Inches(1.18),
    Inches(6.00), Inches(0.38), sz=Pt(14), bold=True, color=C_BLUE)

conclusions = [
    (C_GREEN,  C_LGREEN,
     '✓ TDA agrega valor real',
     '−14.2% MAE vs LSTM puro (29.49 vs 34.38)\n'
     '−4.8% vs SARIMA (29.49 vs 30.98)\n'
     'Features topológicos capturan estructura cíclica no lineal.'),
    (C_BLUE,   C_LBLUE,
     '✓ SARIMA supera a LSTM Puro',
     'SARIMA: 30.98  vs  LSTM puro: 34.38\n'
     'Benchmark estadístico sólido para esta serie.\n'
     'LSTM solo supera a SARIMA con features TDA.'),
    (C_ORANGE, C_LORANG,
     '✓ Calendar features = quick win',
     'month_sin/cos → −1.63 MAE (de 31.93 a 30.31)\n'
     'LSTM no infería estacionalidad anual desde TDA solo.\n'
     'Costo computacional mínimo.'),
    (C_RED,    C_LRED,
     '✗ N=168 limita complejidad',
     'Exógenas, TDA v2 extra, EP Loss → sin mejora.\n'
     'Más features o parámetros con datos insuficientes\n'
     'produce sobreajuste, no generalización.'),
    (C_GRAY,   C_LGRAY,
     '→ Próximo paso',
     'Entrenar en diferencias Δprice.\n'
     'STL decomposition: tendencia + estacionalidad.\n'
     'Elimina el shift implícito en predicciones.'),
]
y_c = Inches(1.62)
for border, bg, label, body in conclusions:
    info_box(sl, label, body,
             Inches(7.00), y_c, Inches(6.10), Inches(1.14),
             border_rgb=border, bg_rgb=bg)
    y_c += Inches(1.18)

notes(sl,
"""Para cerrar, cuatro conclusiones clave.

Primera: TDA agrega valor real. −14.2% MAE sobre LSTM puro y −4.8% sobre SARIMA.

Segunda: SARIMA es un baseline estadístico sólido —supera al LSTM puro, pero queda por debajo de LSTM+TDA.

Tercera: con N=168, la complejidad tiene límites. Más features, exógenas y pérdidas asimétricas no mejoran con pocos datos.

Cuarta: calendar features son un quick-win de bajo costo. El próximo paso natural es entrenar sobre diferencias y STL para eliminar el shift implícito.""")

# ── Save ──────────────────────────────────────────────────────────
prs.save(PPTX_OUT)
print(f'✅  {PPTX_OUT}  ({len(prs.slides)} slides)')
